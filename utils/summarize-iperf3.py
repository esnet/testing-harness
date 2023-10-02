#!/usr/bin/env python3

# This script will summarize the JSON files generated by pscheduler or iperf3.
# For pscheduler, you must run the script "format_pscheduler.sh" first.
# For iperf3 output, use the -i flag

import os
import json
import sys
import argparse
import statistics
import re
from tabulate import tabulate
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Parse command-line arguments
parser = argparse.ArgumentParser(description="Summarize results of pscheduler testing .json files")
parser.add_argument("-3", "--iperf3-json", action="store_true", help="Assume iperf3 output, not pscheduler JSON")
parser.add_argument("-2", "--iperf2", action="store_true", help="Assume iperf2 output from pscheduler")
parser.add_argument("-f", "--full-results", action="store_true", help="show full results for each individual test")
parser.add_argument("-o", "--output-pptx", help="Specify the output PowerPoint file")
args = parser.parse_args()

# add as a flag
VERBOSE = 0
#VERBOSE = 1

# Specify the directory to start the search from
directory_path = "."

# Create a list to store the extracted data
data = []

# Create a nested dictionary to store and calculate the average throughput and retransmits
average_throughput = {}
pacing_string = ""

###########################################################################################
# Define a function to create the PowerPoint presentation
#
def create_pptx(data, avg_headers):
    prs = Presentation()

    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)

    # Create a title slide
    title_slide_layout = prs.slide_layouts[0]
    #blank_slide_layout = prs.slide_layouts[6]
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Performance Test Summary"
    subtitle.text = "Generated by summarize-iperf3-pptx.py"

    # Create a blank slide to start with
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(1)
    top = Inches(1.5)
    # table size (separate from presentation size)
    width = Inches(11)
    height = Inches(6)
    current_dest_host = None

    # Track the number of hosts processed
    hosts_processed = 0

    # Set column widths
    col_widths = [Inches(2.75), Inches(1), Inches(1.25), Inches(1.25), Inches(1.25), Inches(1.25), Inches(1.25), Inches(1.25)]

    # Iterate through data
    # NOTE: this assume 6 test types per host. May need to edit this, or add a flag, or try to compute
    rows = 6  # this many data rows per slide (1 header, 6 data rows)
    slide_num = 1
    for row_idx, row_data in enumerate(data):

        # Check if we need to start a new slide
        dest_host = row_data[0]
        if dest_host != current_dest_host:
            current_dest_host = dest_host
        # old way
        #if hosts_processed == 0 or hosts_processed % rows == 0 : # first time, and after every 6 hosts
            print (f"Creating slide number: {slide_num}, host: {dest_host} ")
            # Create a new slide
            slide = prs.slides.add_slide(blank_slide_layout)
            title = slide.shapes.title
            title.text = dest_host

            current_row = 0
            slide_num += 1
            table = slide.shapes.add_table(rows+1, cols=len(avg_headers), left=left, top=top, width=width, height=height).table
            table.style = "Table Grid"
            for i, width in enumerate(col_widths):
                table.columns[i].width = width
            for col_idx, header_text in enumerate(avg_headers):
                cell = table.cell(0, col_idx)
                cell.text = header_text
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Header cell color
            hosts_processed += 1

        # Populate table data
        tr = table.rows[current_row]
        tr.height = Inches(.5) # set row height
        for col_idx, cell_data in enumerate(row_data):
            #print (f"working on row {current_row + 1} column {col_idx}")
            #print (cell_data)
            cell = table.cell(current_row + 1, col_idx)  # add 1 to not overwrite header
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        current_row += 1

    # Save the PPTX file
    pptx_output_path = args.output_pptx
    prs.save(pptx_output_path)

###########################################################################################
# Traverse all directories and files recursively
for root, dirs, files in os.walk(directory_path):
    if root != "." or dirs == []:
        pacing_string = ""
        if not args.iperf3_json:
           # first check jobmeta.json file for pacing for set of tests in this dir
            file_path = os.path.join(root, "jobmeta.json")
            try:
                with open(file_path, "r") as json_file:
                    json_data = json.load(json_file)
            except:
                print("JSON jobmeta error. Not a JSON file?", file_path)
                continue
            pacing_string = json_data["pacing"]
            if pacing_string:
                pacing_int = int(re.search(r'\d+', pacing_string).group())
#                   print ("Pacing set to: ", pacing_int)
            else:
                pacing_int = 0

    json_file_cnt = 0
    for filename in files:

        file_path = os.path.join(root, filename)
        if VERBOSE:
            print ("Processing file: ", file_path)
        # skip meta files
        if "meta" in filename or "pacing" in filename:
            continue
        if os.path.getsize(file_path) < 1000 or filename.startswith("tmp") or filename.startswith("test"):
            if VERBOSE:
                print(f"   Skipping '{file_path}' because it's less than 1000 bytes.")
            continue 
        if args.iperf2:
            print ("iperf2 file:", file_path)
        elif not args.iperf3_json and not filename.endswith(".json") :
            if VERBOSE:
                print ("    Skipping file: ", file_path)
                continue  # Skip non-.json files

        if VERBOSE:
            print ("Processing file: ", file_path)

        # get RTT from meta file if found
        parts = filename.split(":")
        try:
            host_meta = os.path.join(root,parts[1] + "-meta.json")
        except:
            print ("ERROR: hostname not found in filename: ",filename)
            #sys.exit()
        try:
            with open(host_meta, "r") as json_file:
                json_data = json.load(json_file)
                rtt = json_data["rtt"]
                dest_host = json_data["hostname"]
                #print(f"Got RTT of %d from meta file %s" % (rtt, host_meta))
        except:
            print("Error getting RTT from file: ", host_meta)
            rtt = 0

        # Open and load the JSON file
        with open(file_path, "r") as json_file:
            try:
                json_data = json.load(json_file)
            except:
                print("JSON load error. Not a JSON file?", file_path)
                #sys.exit(-1)
                # just continue on error: sometimes files are corrupt
                continue

        if VERBOSE:
             print ("Getting data from file: ", file_path)
        if args.iperf2:
             # Extract throughput-bits for SUM stream
             sum_stream = next((stream for stream in json_data["summary"]["streams"] if stream["stream-id"] == "SUM"), None)
             if sum_stream:
                 gbits_per_second = sum_stream["throughput-bits"] / 1000000000
                 print("   Throughput Bits for SUM Stream:", gbits_per_second)
             else:
                 print("   SUM Stream not found in the JSON data.")
                 gbits_per_second = 0
             retransmits = 0  # pscheduler does not capture retrans from iperf2
             fq_rate = float(pacing_int)  # from jobmeta.json file

             # Extract the string after "-Z" in "diags"
             diags_string = json_data["diags"]
             match = re.search(r'-Z (\S+)', diags_string)
             if match:
                 cong = match.group(1)
                 print("   Congestion Control Algorithm:", cong)
             else:
                 print("   Congestion Control Algorithm not found in the 'diags' field.")
             # Extract the string after "-P" in "diags"
             nstreams = json_data["diags"]
             match = re.search(r'-P (\d+)', diags_string)
             if match:
                 nstreams = match.group(1)
                 #print("   Number of Streams:", nstreams)
             else:
                 print("   Number of streams not found in the 'diags' field.")
        else:
            try:
               nstreams = json_data["start"]["test_start"]["num_streams"]
               fq_rate = float(json_data["start"]["test_start"]["fqrate"]) / 1000000000
               cong = json_data["end"]["sender_tcp_congestion"]
            except:
               # test most have failed
               if VERBOSE:
                   print ("Error extracting dest_host from JSON file: ", file_path)
               continue

            gbits_per_second = float(json_data["end"]["sum_sent"]["bits_per_second"]) / 1000000000
            retransmits = json_data["end"]["sum_sent"]["retransmits"]

            # Create a key for the nested dictionary based on dest_host, nstreams, cong, and fq_rate
        if pacing_string:  # pscheduler still does not support fq_rate, so grab it from the testing_harness jobmeta file
            fq_rate = float(pacing_int)
#               print (" setting fq_rate based on jobmeta pacing rate: ",fq_rate)

        key = (dest_host, nstreams, cong, fq_rate)

        # Add the throughput and retransmits to the nested dictionary
        if key in average_throughput:
            average_throughput[key][0].append(gbits_per_second)
            average_throughput[key][1].append(retransmits)
        else:
            average_throughput[key] = [[gbits_per_second], [retransmits], [rtt]]

        # Append the extracted data to the list
        data.append([dest_host, rtt, nstreams, cong, fq_rate, gbits_per_second, retransmits])
        json_file_cnt += 1

if json_file_cnt == 0:
    print ("No JSON data files found. Did you run format_pscheduler.sh? Or do you need the -i flag? ")
    sys.exit()

# Create headers for the average throughput table
avg_headers = ["Dest Host", "RTT (ms)", "Streams", "CC Alg", "Pacing (Gbps)", "Tput (Gbps)", "Stddev (nvals)", "RXMTs"]

# Calculate and format the averages and standard deviation for each combination of dest_host, nstreams, cong, fq_rate
average_table_data = []

for key, values in average_throughput.items():
    avg_throughput = sum(values[0]) / len(values[0])
    avg_retransmits = sum(values[1]) / len(values[1])
    throughput_values = values[0]  # List of throughput values
    data_points = len(throughput_values)  # Number of data points
    if data_points > 1:
        std_dev_throughput = statistics.stdev(throughput_values)  # Calculate stddev
    else:
        std_dev_throughput = 0
    dest_host, nstreams, cong, fq_rate = key
    avg_throughput_formatted = "{:.2f}".format(avg_throughput)
    std_dev_throughput_formatted = "{:.2f}".format(std_dev_throughput)  # Format stddev
    #average_table_data.append([dest_host, nstreams, cong, fq_rate, avg_throughput_formatted, std_dev_throughput_formatted, int(avg_retransmits)])
    rtt = values[2][0]
    average_table_data.append([dest_host, rtt, nstreams, cong, fq_rate, avg_throughput_formatted, std_dev_throughput_formatted + " ({})".format(data_points), int(avg_retransmits)])

# Sort the data by columns: dest_host, nstreams, pacing, cong
average_table_data = sorted(average_table_data, key=lambda x: (x[0], x[2], x[4], x[3]))

# Create a list to store the final formatted data
formatted_average_table_data = []

# Check if the PPTX output option is provided
if args.output_pptx:
    create_pptx(average_table_data, avg_headers)
    print(f"PPTX file '{args.output_pptx}' has been generated.")

# Iterate through the average_table_data and add separators
prev_dest_host = None
for row in average_table_data:
    dest_host, rtt, nstreams, cong, fq_rate, avg_throughput, stddev_throughput, avg_retransmits = row

    # Check if dest_host is different from the previous row
    if dest_host != prev_dest_host:
        # Add a separator row with double lines
        separator_row = ["=" * len(header) for header in avg_headers]
        formatted_average_table_data.extend([separator_row, row])
        prev_dest_host = dest_host  # Update the previous dest_host
    else:
        formatted_average_table_data.append(row)

# Print the average throughput table with double lines for separators
print("\nAverage Throughput for 5 tests:")
print(tabulate(formatted_average_table_data, headers=avg_headers, tablefmt="grid", disable_numparse=True))

# Create headers for the main data table
data_headers = ["Dest Host", "Num Streams", "CC Alg", "Pacing (Gbps)", "Throughput", "Retransmits"]

# Sort the data by columns: dest_host, nstreams, cong, and fq_rate
data = sorted(data, key=lambda x: (x[0], x[1], x[3], x[2]))

if args.full_results:
    # Print the main data table
    print("\nResults from each individual test:")
    print(tabulate(data, headers=data_headers, tablefmt="grid"))


print ("\nDone. \n")
